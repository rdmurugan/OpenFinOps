"""
Enterprise Export & Reporting
=============================

Advanced export capabilities for enterprise visualization including
compliance reporting, branded outputs, and automated distribution.
"""

from __future__ import annotations

import json
import warnings
from dataclasses import dataclass, field
from datetime import datetime
from pathlib import Path
from typing import Any, Dict, List, Optional, Union
import tempfile
import zipfile

try:
    import matplotlib.pyplot as plt
    import matplotlib.backends.backend_pdf as backend_pdf
    MATPLOTLIB_AVAILABLE = True
except ImportError:
    MATPLOTLIB_AVAILABLE = False
    warnings.warn("Matplotlib not available for enterprise exports")

try:
    from reportlab.pdfgen import canvas
    from reportlab.lib.pagesizes import letter, A4
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Image
    from reportlab.lib.styles import getSampleStyleSheet
    REPORTLAB_AVAILABLE = True
except ImportError:
    REPORTLAB_AVAILABLE = False
    warnings.warn("ReportLab not available for PDF reports")

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
    from pptx.dml.color import RGBColor
    PYTHON_PPTX_AVAILABLE = True
except ImportError:
    PYTHON_PPTX_AVAILABLE = False
    warnings.warn("python-pptx not available for PowerPoint exports")

try:
    import openpyxl
    from openpyxl.drawing.image import Image as XLImage
    from openpyxl.styles import Font, Alignment, PatternFill
    OPENPYXL_AVAILABLE = True
except ImportError:
    OPENPYXL_AVAILABLE = False
    warnings.warn("openpyxl not available for Excel exports")

from .charts import ChartMetadata
from .security import SecurityLevel
from .themes import BrandingConfig


@dataclass
class ExportConfig:
    """Configuration for enterprise exports."""
    format: str = "pdf"  # pdf, png, svg, html, json, pptx, xlsx
    dpi: int = 300
    include_metadata: bool = True
    include_audit_trail: bool = True
    watermark: bool = True
    branded: bool = True
    filename_template: str = "{title}_{timestamp}"
    compression: bool = False
    slide_layout: str = "title_and_content"  # For PowerPoint
    excel_worksheet: str = "Charts"  # For Excel


@dataclass
class ReportSection:
    """Section of an enterprise report."""
    title: str
    content_type: str  # chart, text, table, image
    content: Any
    description: Optional[str] = None
    security_level: SecurityLevel = SecurityLevel.INTERNAL


class EnterpriseExporter:
    """
    Enterprise-grade export system with compliance features,
    branding, and audit trails.
    """

    def __init__(self, branding: Optional[BrandingConfig] = None):
        self.branding = branding or BrandingConfig()
        self.temp_dir = Path(tempfile.mkdtemp())

    def export_chart(self, figure, metadata: ChartMetadata,
                    config: ExportConfig) -> str:
        """
        Export single chart with enterprise features.

        Returns path to exported file.
        """
        if not MATPLOTLIB_AVAILABLE:
            raise ImportError("Matplotlib required for chart export")

        # Generate filename
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        filename = config.filename_template.format(
            title=metadata.title.replace(" ", "_"),
            timestamp=timestamp
        )

        output_path = self.temp_dir / f"{filename}.{config.format}"

        if config.format in ['png', 'jpg', 'svg']:
            self._export_image(figure, output_path, config, metadata)
        elif config.format == 'pdf':
            self._export_pdf_chart(figure, output_path, config, metadata)
        elif config.format == 'html':
            self._export_html_chart(figure, output_path, config, metadata)
        elif config.format == 'json':
            self._export_json_metadata(output_path, metadata)
        elif config.format == 'pptx':
            self._export_powerpoint_chart(figure, output_path, config, metadata)
        elif config.format == 'xlsx':
            self._export_excel_chart(figure, output_path, config, metadata)
        else:
            raise ValueError(f"Unsupported export format: {config.format}")

        return str(output_path)

    def _export_image(self, figure, output_path: Path, config: ExportConfig,
                     metadata: ChartMetadata) -> None:
        """Export chart as image with enterprise features."""
        # Add branding if requested
        if config.branded:
            self._add_enterprise_branding(figure)

        # Add watermark if security level requires it
        if config.watermark and metadata.security_level in [
            SecurityLevel.CONFIDENTIAL, SecurityLevel.RESTRICTED
        ]:
            self._add_security_watermark(figure, metadata.security_level)

        # Handle VizlyFigure vs matplotlib Figure
        fig = figure.figure if hasattr(figure, 'figure') else figure

        # Save with high DPI
        fig.savefig(
            output_path,
            dpi=config.dpi,
            bbox_inches='tight',
            facecolor='white',
            edgecolor='none'
        )

        # Add metadata to file if supported
        if config.include_metadata and output_path.suffix.lower() == '.png':
            self._embed_png_metadata(output_path, metadata)

    def _export_pdf_chart(self, figure, output_path: Path, config: ExportConfig,
                         metadata: ChartMetadata) -> None:
        """Export chart as PDF with compliance features."""
        if config.branded:
            self._add_enterprise_branding(figure)

        # Create PDF with metadata
        with backend_pdf.PdfPages(output_path, metadata={
            'Title': metadata.title,
            'Author': metadata.created_by,
            'Subject': f"Enterprise Chart - {metadata.title}",
            'Creator': f"Vizly Enterprise - {self.branding.company_name}",
            'CreationDate': metadata.created_at,
            'ModDate': metadata.last_modified
        }) as pdf:
            # Add chart page
            pdf.savefig(figure, bbox_inches='tight')

            # Add metadata page if requested
            if config.include_metadata:
                self._add_metadata_page(pdf, metadata)

    def _export_html_chart(self, figure, output_path: Path, config: ExportConfig,
                          metadata: ChartMetadata) -> None:
        """Export chart as interactive HTML."""
        # Save chart as base64 image
        import base64
        import io

        # Handle VizlyFigure vs matplotlib Figure
        fig = figure.figure if hasattr(figure, 'figure') else figure

        img_buffer = io.BytesIO()
        fig.savefig(img_buffer, format='png', dpi=config.dpi, bbox_inches='tight')
        img_str = base64.b64encode(img_buffer.getvalue()).decode()

        # Generate HTML
        html_content = self._generate_html_report([
            ReportSection(
                title=metadata.title,
                content_type="chart",
                content=img_str,
                description=f"Created by {metadata.created_by} on {metadata.created_at.strftime('%Y-%m-%d')}"
            )
        ], metadata, config)

        with open(output_path, 'w', encoding='utf-8') as f:
            f.write(html_content)

    def _export_json_metadata(self, output_path: Path, metadata: ChartMetadata) -> None:
        """Export chart metadata as JSON."""
        metadata_dict = {
            'chart_id': metadata.chart_id,
            'title': metadata.title,
            'created_by': metadata.created_by,
            'created_at': metadata.created_at.isoformat(),
            'last_modified': metadata.last_modified.isoformat(),
            'security_level': metadata.security_level.value,
            'data_sources': metadata.data_sources,
            'compliance_tags': metadata.compliance_tags,
            'approved_by': metadata.approved_by,
            'approval_date': metadata.approval_date.isoformat() if metadata.approval_date else None
        }

        with open(output_path, 'w', encoding='utf-8') as f:
            json.dump(metadata_dict, f, indent=2)

    def create_executive_report(self, sections: List[ReportSection],
                              title: str = "Executive Report",
                              format: str = "pdf") -> str:
        """Create multi-section executive report."""
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        filename = f"executive_report_{timestamp}"

        if format == "pdf" and REPORTLAB_AVAILABLE:
            return self._create_pdf_report(sections, title, filename)
        elif format == "html":
            return self._create_html_report(sections, title, filename)
        else:
            raise ValueError(f"Unsupported report format: {format}")

    def _create_pdf_report(self, sections: List[ReportSection],
                          title: str, filename: str) -> str:
        """Create comprehensive PDF report."""
        if not REPORTLAB_AVAILABLE:
            raise ImportError("ReportLab required for PDF reports")

        output_path = self.temp_dir / f"{filename}.pdf"
        doc = SimpleDocTemplate(str(output_path), pagesize=A4)
        styles = getSampleStyleSheet()
        story = []

        # Title page
        story.append(Paragraph(title, styles['Title']))
        story.append(Spacer(1, 20))
        story.append(Paragraph(f"Generated by {self.branding.company_name}", styles['Normal']))
        story.append(Paragraph(f"Date: {datetime.now().strftime('%Y-%m-%d %H:%M')}", styles['Normal']))
        story.append(Spacer(1, 40))

        # Table of contents
        story.append(Paragraph("Table of Contents", styles['Heading1']))
        for i, section in enumerate(sections, 1):
            story.append(Paragraph(f"{i}. {section.title}", styles['Normal']))
        story.append(Spacer(1, 30))

        # Sections
        for i, section in enumerate(sections, 1):
            story.append(Paragraph(f"{i}. {section.title}", styles['Heading1']))

            if section.description:
                story.append(Paragraph(section.description, styles['Normal']))
                story.append(Spacer(1, 12))

            if section.content_type == "chart":
                # Add chart image
                chart_path = self.temp_dir / f"temp_chart_{i}.png"
                section.content.savefig(chart_path, dpi=300, bbox_inches='tight')
                story.append(Image(str(chart_path), width=400, height=300))

            elif section.content_type == "text":
                story.append(Paragraph(section.content, styles['Normal']))

            story.append(Spacer(1, 20))

        # Build PDF
        doc.build(story)
        return str(output_path)

    def _create_html_report(self, sections: List[ReportSection],
                           title: str, filename: str) -> str:
        """Create HTML report with interactive features."""
        output_path = self.temp_dir / f"{filename}.html"

        # Generate metadata for HTML
        dummy_metadata = ChartMetadata(
            chart_id="report",
            title=title,
            created_by="system"
        )

        config = ExportConfig(format="html")
        html_content = self._generate_html_report(sections, dummy_metadata, config)

        with open(output_path, 'w', encoding='utf-8') as f:
            f.write(html_content)

        return str(output_path)

    def _generate_html_report(self, sections: List[ReportSection],
                             metadata: ChartMetadata, config: ExportConfig) -> str:
        """Generate HTML content for reports."""
        html = f"""
<!DOCTYPE html>
<html lang="en">
<head>
    <meta charset="UTF-8">
    <meta name="viewport" content="width=device-width, initial-scale=1.0">
    <title>{metadata.title} - {self.branding.company_name}</title>
    <style>
        body {{
            font-family: {self.branding.font_family}, sans-serif;
            margin: 0;
            padding: 20px;
            background-color: #f8f9fa;
        }}
        .container {{
            max-width: 1200px;
            margin: 0 auto;
            background: white;
            padding: 30px;
            border-radius: 8px;
            box-shadow: 0 2px 10px rgba(0,0,0,0.1);
        }}
        .header {{
            border-bottom: 3px solid {self.branding.primary_color};
            padding-bottom: 20px;
            margin-bottom: 30px;
        }}
        .title {{
            color: {self.branding.primary_color};
            margin: 0;
            font-size: 2.5em;
        }}
        .subtitle {{
            color: #666;
            margin: 10px 0 0 0;
        }}
        .section {{
            margin: 40px 0;
            padding: 20px;
            border-left: 4px solid {self.branding.accent_color};
            background: #fafafa;
        }}
        .section-title {{
            color: {self.branding.primary_color};
            font-size: 1.5em;
            margin-bottom: 15px;
        }}
        .chart-container {{
            text-align: center;
            margin: 20px 0;
        }}
        .chart-image {{
            max-width: 100%;
            border: 1px solid #ddd;
            border-radius: 4px;
        }}
        .metadata {{
            background: #e9ecef;
            padding: 15px;
            border-radius: 4px;
            font-size: 0.9em;
            margin-top: 30px;
        }}
        .security-banner {{
            background: #fff3cd;
            border: 1px solid #ffeaa7;
            color: #856404;
            padding: 10px;
            border-radius: 4px;
            margin-bottom: 20px;
            font-weight: bold;
        }}
        @media print {{
            body {{ background: white; }}
            .container {{ box-shadow: none; }}
        }}
    </style>
</head>
<body>
    <div class="container">
        <div class="header">
            <h1 class="title">{metadata.title}</h1>
            <p class="subtitle">Generated by {self.branding.company_name} • {datetime.now().strftime('%Y-%m-%d %H:%M')}</p>
        </div>
"""

        # Security banner if needed
        if metadata.security_level in [SecurityLevel.CONFIDENTIAL, SecurityLevel.RESTRICTED]:
            html += f"""
        <div class="security-banner">
            🔒 {metadata.security_level.value.upper()} - Authorized Personnel Only
        </div>
"""

        # Add sections
        for section in sections:
            html += f"""
        <div class="section">
            <h2 class="section-title">{section.title}</h2>
"""
            if section.description:
                html += f"<p>{section.description}</p>"

            if section.content_type == "chart":
                html += f"""
            <div class="chart-container">
                <img src="data:image/png;base64,{section.content}" class="chart-image" alt="{section.title}">
            </div>
"""
            elif section.content_type == "text":
                html += f"<p>{section.content}</p>"

            html += "</div>"

        # Add metadata if requested
        if config.include_metadata:
            html += f"""
        <div class="metadata">
            <h3>Document Metadata</h3>
            <p><strong>Chart ID:</strong> {metadata.chart_id}</p>
            <p><strong>Created by:</strong> {metadata.created_by}</p>
            <p><strong>Created:</strong> {metadata.created_at.strftime('%Y-%m-%d %H:%M')}</p>
            <p><strong>Security Level:</strong> {metadata.security_level.value}</p>
            <p><strong>Data Sources:</strong> {', '.join(metadata.data_sources) if metadata.data_sources else 'None specified'}</p>
            <p><strong>Compliance Tags:</strong> {', '.join(metadata.compliance_tags) if metadata.compliance_tags else 'None'}</p>
        </div>
"""

        html += """
    </div>
</body>
</html>"""

        return html

    def _add_enterprise_branding(self, figure) -> None:
        """Add enterprise branding elements to figure."""
        # Handle VizlyFigure vs matplotlib Figure
        fig = figure.figure if hasattr(figure, 'figure') else figure

        # Add company name in footer
        fig.text(0.02, 0.02, f"© {self.branding.company_name}",
                fontsize=8, alpha=0.7, color='gray')

        # Add generation timestamp
        timestamp = datetime.now().strftime('%Y-%m-%d %H:%M')
        fig.text(0.98, 0.02, f"Generated: {timestamp}",
                ha='right', fontsize=8, alpha=0.7, color='gray')

    def _add_security_watermark(self, figure, security_level: SecurityLevel) -> None:
        """Add security classification watermark."""
        watermark_text = f"{security_level.value.upper()} - AUTHORIZED PERSONNEL ONLY"

        figure.text(0.5, 0.5, watermark_text, ha='center', va='center',
                   fontsize=20, alpha=0.1, rotation=45,
                   transform=figure.transFigure, color='red')

    def _embed_png_metadata(self, png_path: Path, metadata: ChartMetadata) -> None:
        """Embed metadata in PNG file."""
        try:
            from PIL import Image
            from PIL.PngImagePlugin import PngInfo

            # Load image
            img = Image.open(png_path)

            # Create metadata
            png_info = PngInfo()
            png_info.add_text("Title", metadata.title)
            png_info.add_text("Author", metadata.created_by)
            png_info.add_text("Description", f"Enterprise chart - {metadata.title}")
            png_info.add_text("Software", f"Vizly Enterprise - {self.branding.company_name}")
            png_info.add_text("Security", metadata.security_level.value)

            # Save with metadata
            img.save(png_path, "PNG", pnginfo=png_info)

        except ImportError:
            # PIL not available, skip metadata embedding
            pass

    def _add_metadata_page(self, pdf, metadata: ChartMetadata) -> None:
        """Add metadata page to PDF."""
        # Create new figure for metadata
        fig, ax = plt.subplots(figsize=(8, 10))
        ax.axis('off')

        # Add metadata text
        metadata_text = f"""
Chart Metadata

Title: {metadata.title}
Chart ID: {metadata.chart_id}
Created by: {metadata.created_by}
Created: {metadata.created_at.strftime('%Y-%m-%d %H:%M:%S')}
Last Modified: {metadata.last_modified.strftime('%Y-%m-%d %H:%M:%S')}
Security Level: {metadata.security_level.value}

Data Sources:
{chr(10).join(f"• {source}" for source in metadata.data_sources) if metadata.data_sources else "None specified"}

Compliance Tags:
{chr(10).join(f"• {tag}" for tag in metadata.compliance_tags) if metadata.compliance_tags else "None"}

Approval:
Approved by: {metadata.approved_by or "Pending"}
Approval date: {metadata.approval_date.strftime('%Y-%m-%d') if metadata.approval_date else "Pending"}

Generated by Vizly Enterprise
© {self.branding.company_name}
        """

        ax.text(0.1, 0.9, metadata_text, transform=ax.transAxes,
               fontsize=10, verticalalignment='top', fontfamily='monospace')

        pdf.savefig(fig, bbox_inches='tight')
        plt.close(fig)

    def create_compliance_package(self, charts: List[Tuple[Any, ChartMetadata]],
                                 package_name: str = "compliance_package") -> str:
        """Create compliance package with all charts and documentation."""
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        package_path = self.temp_dir / f"{package_name}_{timestamp}.zip"

        with zipfile.ZipFile(package_path, 'w', zipfile.ZIP_DEFLATED) as zip_file:
            # Add charts
            for i, (figure, metadata) in enumerate(charts):
                config = ExportConfig(format="pdf", include_metadata=True)
                chart_path = self.export_chart(figure, metadata, config)
                zip_file.write(chart_path, f"charts/chart_{i+1:03d}_{metadata.title}.pdf")

                # Add metadata JSON
                json_config = ExportConfig(format="json")
                json_path = self.export_chart(figure, metadata, json_config)
                zip_file.write(json_path, f"metadata/chart_{i+1:03d}_metadata.json")

            # Add compliance summary
            summary = self._generate_compliance_summary(charts)
            summary_path = self.temp_dir / "compliance_summary.json"
            with open(summary_path, 'w') as f:
                json.dump(summary, f, indent=2)
            zip_file.write(summary_path, "compliance_summary.json")

        return str(package_path)

    def _export_powerpoint_chart(self, figure, output_path: Path, config: ExportConfig,
                                metadata: ChartMetadata) -> None:
        """Export chart to PowerPoint presentation."""
        if not PYTHON_PPTX_AVAILABLE:
            raise ImportError("python-pptx required for PowerPoint export")

        # Create presentation with branded template
        prs = Presentation()

        # Apply branding theme
        self._apply_powerpoint_branding(prs)

        # Create slide
        if config.slide_layout == "title_and_content":
            slide_layout = prs.slide_layouts[1]  # Title and Content layout
        elif config.slide_layout == "blank":
            slide_layout = prs.slide_layouts[6]  # Blank layout
        else:
            slide_layout = prs.slide_layouts[0]  # Title slide

        slide = prs.slides.add_slide(slide_layout)

        # Add title
        if hasattr(slide.shapes, 'title') and slide.shapes.title:
            slide.shapes.title.text = metadata.title

            # Apply branding to title
            title_font = slide.shapes.title.text_frame.paragraphs[0].font
            title_font.name = self.branding.font_family
            title_font.size = Pt(24)
            title_font.color.rgb = RGBColor.from_string(self.branding.primary_color.lstrip('#'))

        # Save chart as image first
        temp_chart_path = self.temp_dir / "temp_chart.png"
        if config.branded:
            self._add_enterprise_branding(figure)

        # Handle VizlyFigure vs matplotlib Figure
        fig = figure.figure if hasattr(figure, 'figure') else figure
        fig.savefig(temp_chart_path, dpi=config.dpi, bbox_inches='tight',
                   facecolor='white', edgecolor='none')

        # Add chart image to slide
        left = Inches(1)
        top = Inches(2) if hasattr(slide.shapes, 'title') and slide.shapes.title else Inches(0.5)
        height = Inches(5)

        # Calculate width maintaining aspect ratio
        from PIL import Image as PILImage
        with PILImage.open(temp_chart_path) as img:
            aspect_ratio = img.width / img.height
            width = height * aspect_ratio

        slide.shapes.add_picture(str(temp_chart_path), left, top, width, height)

        # Add metadata text box if requested
        if config.include_metadata:
            metadata_text = (
                f"Created by: {metadata.created_by}\n"
                f"Date: {metadata.created_at.strftime('%Y-%m-%d')}\n"
                f"Security: {metadata.security_level.value}"
            )

            # Add text box in bottom right
            left = Inches(6)
            top = Inches(7)
            width = Inches(3)
            height = Inches(1)

            textbox = slide.shapes.add_textbox(left, top, width, height)
            text_frame = textbox.text_frame
            text_frame.text = metadata_text

            # Format metadata text
            for paragraph in text_frame.paragraphs:
                paragraph.font.size = Pt(9)
                paragraph.font.color.rgb = RGBColor(128, 128, 128)

        # Add security watermark if needed
        if config.watermark and metadata.security_level in [
            SecurityLevel.CONFIDENTIAL, SecurityLevel.RESTRICTED
        ]:
            self._add_powerpoint_watermark(slide, metadata.security_level)

        # Save presentation
        prs.save(output_path)

    def _export_excel_chart(self, figure, output_path: Path, config: ExportConfig,
                           metadata: ChartMetadata) -> None:
        """Export chart to Excel workbook."""
        if not OPENPYXL_AVAILABLE:
            raise ImportError("openpyxl required for Excel export")

        # Create workbook
        wb = openpyxl.Workbook()
        ws = wb.active
        ws.title = config.excel_worksheet

        # Apply branding theme
        self._apply_excel_branding(ws)

        # Add title
        ws['A1'] = metadata.title
        title_font = Font(name=self.branding.font_family, size=16, bold=True,
                         color=self.branding.primary_color.lstrip('#'))
        ws['A1'].font = title_font
        ws['A1'].alignment = Alignment(horizontal='center')

        # Merge title cells
        ws.merge_cells('A1:F1')

        # Save chart as image
        temp_chart_path = self.temp_dir / "temp_excel_chart.png"
        if config.branded:
            self._add_enterprise_branding(figure)

        # Handle VizlyFigure vs matplotlib Figure
        fig = figure.figure if hasattr(figure, 'figure') else figure
        fig.savefig(temp_chart_path, dpi=config.dpi, bbox_inches='tight',
                   facecolor='white', edgecolor='none')

        # Insert image
        img = XLImage(str(temp_chart_path))
        # Scale image to fit well in Excel
        img.height = 300  # pixels
        img.width = int(img.width * (300 / img.height))  # maintain aspect ratio

        ws.add_image(img, 'A3')

        # Add metadata table
        if config.include_metadata:
            metadata_row = 20  # Below the chart

            # Metadata headers
            headers = ['Property', 'Value']
            for col, header in enumerate(headers, 1):
                cell = ws.cell(row=metadata_row, column=col)
                cell.value = header
                cell.font = Font(bold=True)
                cell.fill = PatternFill(start_color='E6E6FA', end_color='E6E6FA', fill_type='solid')

            # Metadata values
            metadata_items = [
                ('Chart ID', metadata.chart_id),
                ('Created By', metadata.created_by),
                ('Created Date', metadata.created_at.strftime('%Y-%m-%d %H:%M')),
                ('Last Modified', metadata.last_modified.strftime('%Y-%m-%d %H:%M')),
                ('Security Level', metadata.security_level.value),
                ('Data Sources', ', '.join(metadata.data_sources) if metadata.data_sources else 'None'),
                ('Compliance Tags', ', '.join(metadata.compliance_tags) if metadata.compliance_tags else 'None')
            ]

            for row, (prop, value) in enumerate(metadata_items, metadata_row + 1):
                ws.cell(row=row, column=1).value = prop
                ws.cell(row=row, column=2).value = value

            # Auto-adjust column widths
            for column in ws.columns:
                max_length = 0
                column_letter = column[0].column_letter
                for cell in column:
                    try:
                        if len(str(cell.value)) > max_length:
                            max_length = len(str(cell.value))
                    except:
                        pass
                adjusted_width = min(max_length + 2, 50)
                ws.column_dimensions[column_letter].width = adjusted_width

        # Add company footer
        footer_text = f"© {self.branding.company_name} - Generated on {datetime.now().strftime('%Y-%m-%d %H:%M')}"
        last_row = ws.max_row + 2
        ws.cell(row=last_row, column=1).value = footer_text
        footer_font = Font(size=8, color='808080')
        ws.cell(row=last_row, column=1).font = footer_font

        # Save workbook
        wb.save(output_path)

    def _apply_powerpoint_branding(self, presentation) -> None:
        """Apply enterprise branding to PowerPoint presentation."""
        # Set presentation properties
        core_props = presentation.core_properties
        core_props.author = self.branding.company_name
        core_props.title = f"Enterprise Report - {self.branding.company_name}"
        core_props.subject = "Generated by Vizly Enterprise"

    def _apply_excel_branding(self, worksheet) -> None:
        """Apply enterprise branding to Excel worksheet."""
        # Set worksheet properties via workbook
        wb = worksheet.parent
        wb.properties.creator = self.branding.company_name
        wb.properties.title = f"Enterprise Charts - {self.branding.company_name}"
        wb.properties.subject = "Generated by Vizly Enterprise"

    def _add_powerpoint_watermark(self, slide, security_level: SecurityLevel) -> None:
        """Add security watermark to PowerPoint slide."""
        watermark_text = f"{security_level.value.upper()} - AUTHORIZED PERSONNEL ONLY"

        # Add diagonal watermark text
        left = Inches(2)
        top = Inches(3)
        width = Inches(6)
        height = Inches(1)

        textbox = slide.shapes.add_textbox(left, top, width, height)
        text_frame = textbox.text_frame
        text_frame.text = watermark_text

        # Format watermark
        paragraph = text_frame.paragraphs[0]
        paragraph.font.size = Pt(48)
        paragraph.font.color.rgb = RGBColor(255, 0, 0)  # Red
        paragraph.alignment = PP_ALIGN.CENTER

        # Make semi-transparent and rotated (approximation since python-pptx has limited rotation)
        textbox.fill.solid()
        textbox.fill.fore_color.rgb = RGBColor(255, 255, 255)
        textbox.fill.transparency = 0.7

    def create_powerpoint_report(self, sections: List[ReportSection],
                                title: str = "Executive Report") -> str:
        """Create multi-slide PowerPoint report."""
        if not PYTHON_PPTX_AVAILABLE:
            raise ImportError("python-pptx required for PowerPoint reports")

        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        filename = f"powerpoint_report_{timestamp}"
        output_path = self.temp_dir / f"{filename}.pptx"

        # Create presentation
        prs = Presentation()
        self._apply_powerpoint_branding(prs)

        # Title slide
        title_slide_layout = prs.slide_layouts[0]
        title_slide = prs.slides.add_slide(title_slide_layout)
        title_slide.shapes.title.text = title
        subtitle = title_slide.placeholders[1]
        subtitle.text = f"{self.branding.company_name}\n{datetime.now().strftime('%B %d, %Y')}"

        # Content slides
        for i, section in enumerate(sections):
            slide_layout = prs.slide_layouts[1]  # Title and Content
            slide = prs.slides.add_slide(slide_layout)

            # Title
            slide.shapes.title.text = section.title

            if section.content_type == "chart":
                # Save chart as temp image
                temp_path = self.temp_dir / f"temp_slide_{i}.png"
                if hasattr(section.content, 'savefig'):
                    section.content.savefig(temp_path, dpi=300, bbox_inches='tight')
                elif hasattr(section.content, 'save'):
                    section.content.save(temp_path)

                # Add image to slide
                left = Inches(1)
                top = Inches(2)
                height = Inches(5)
                slide.shapes.add_picture(str(temp_path), left, top, height=height)

            elif section.content_type == "text":
                # Add text content
                if len(slide.placeholders) > 1:
                    content_placeholder = slide.placeholders[1]
                    content_placeholder.text = section.content

        # Save presentation
        prs.save(output_path)
        return str(output_path)

    def _generate_compliance_summary(self, charts: List[Tuple[Any, ChartMetadata]]) -> Dict[str, Any]:
        """Generate compliance summary for package."""
        return {
            "package_generated": datetime.now().isoformat(),
            "total_charts": len(charts),
            "security_levels": {
                level.value: sum(1 for _, metadata in charts
                               if metadata.security_level == level)
                for level in SecurityLevel
            },
            "compliance_tags": list(set(
                tag for _, metadata in charts
                for tag in metadata.compliance_tags
            )),
            "data_sources": list(set(
                source for _, metadata in charts
                for source in metadata.data_sources
            )),
            "approval_status": {
                "approved": sum(1 for _, metadata in charts if metadata.approved_by),
                "pending": sum(1 for _, metadata in charts if not metadata.approved_by)
            }
        }